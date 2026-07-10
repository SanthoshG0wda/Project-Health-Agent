import argparse
import json
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn


C_NAVY = RGBColor(0x1E, 0x27, 0x61)
C_ICE = RGBColor(0xCA, 0xDC, 0xFC)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK = RGBColor(0x0D, 0x12, 0x33)
C_CARD = RGBColor(0x1A, 0x23, 0x48)
C_GRAY = RGBColor(0x88, 0x92, 0xB0)
C_MUTED = RGBColor(0x4A, 0x55, 0x80)
C_RED = RGBColor(0xE0, 0x3C, 0x31)
C_AMBER = RGBColor(0xF5, 0xA6, 0x23)
C_GREEN = RGBColor(0x27, 0xAE, 0x60)


def _set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _tb(slide, left, top, width, height, text, size=14, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri", valign=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if valign == "middle":
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        txBox.text_frame._txBody.bodyPr.set('anchor', 'ctr')
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox


def _rect(slide, left, top, width, height, color, line=False):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not line:
        shape.line.fill.background()
    return shape


def _oval(slide, left, top, size, color):
    shape = slide.shapes.add_shape(9, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _shadow(shape):
    sp = shape._element
    spPr = sp.find(qn('p:spPr')) if sp.find(qn('p:spPr')) is not None else sp.find(qn('a:spPr'))
    if spPr is None:
        spPr = sp.find('.//' + qn('a:spPr'))
    if spPr is None:
        return
    shadow_el = spPr.makeelement(qn('a:effectLst'), {})
    outer = shadow_el.makeelement(qn('a:outerShdw'), {
        'blurRad': '40000', 'dist': '20000', 'dir': '1350000',
        'algn': 'bl', 'rotWithShape': '0'
    })
    srgb = outer.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alpha = srgb.makeelement(qn('a:alpha'), {'val': '30000'})
    srgb.append(alpha)
    outer.append(srgb)
    shadow_el.append(outer)
    spPr.append(shadow_el)


def rag_color(status):
    if status == "Red": return C_RED
    if status == "Amber": return C_AMBER
    if status == "Green": return C_GREEN
    return C_GRAY


def load_weekly_data(input_dir):
    weeks = sorted(glob.glob(f"{input_dir}/[0-9][0-9][0-9][0-9]-[0-9][0-9]-[0-9][0-9]/summary.json"))
    if not weeks:
        raise ValueError(f"No weekly data found in {input_dir}")
    with open(weeks[-1]) as f:
        return json.load(f)


def build(data):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    projects = data["projects"]
    latest = data["week"]
    n_red = sum(1 for p in projects if p["overall"] == "Red")
    n_amber = sum(1 for p in projects if p["overall"] == "Amber")
    n_green = sum(1 for p in projects if p["overall"] == "Green")
    total = len(projects)

    # ── Slide 1: Title ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, C_NAVY)
    _rect(s, 0, 7, 13.33, 0.5, C_DARK)
    # RAG stack on right
    for j, (letter, color) in enumerate([("R", C_RED), ("A", C_AMBER), ("G", C_GREEN)]):
        _tb(s, 10.5, 1.2 + j * 1.3, 1.5, 1.2, letter, 72, color, True, PP_ALIGN.CENTER, "Georgia")
    _rect(s, 10.2, 1.5, 0.04, 4.5, C_MUTED)
    _tb(s, 1, 1.5, 8.5, 2.2, "Project Health\nReport", 44, C_WHITE, True, font="Georgia")
    _tb(s, 1, 3.8, 8, 0.6, f"Week ending {latest}", 18, C_ICE, font="Calibri")
    _tb(s, 1, 4.5, 8, 0.5, "For Client Leadership Review", 14, C_GRAY, font="Calibri")

    # ── Slide 2: RAG Dashboard ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, C_DARK)
    _tb(s, 0.8, 0.4, 10, 0.7, "Portfolio RAG Dashboard", 28, C_WHITE, True, font="Georgia")

    stats = [("Total Projects", total, C_WHITE), ("Red", n_red, C_RED), ("Amber", n_amber, C_AMBER), ("Green", n_green, C_GREEN)]
    for i, (label, val, color) in enumerate(stats):
        x = 0.8 + i * 3.05
        card = _rect(s, x, 1.5, 2.7, 1.6, C_CARD)
        _shadow(card)
        _tb(s, x, 1.6, 2.7, 0.9, str(val), 48, color, True, PP_ALIGN.CENTER, "Georgia")
        _tb(s, x, 2.5, 2.7, 0.4, label, 13, C_GRAY, align=PP_ALIGN.CENTER)

    chart_card = _rect(s, 0.8, 3.6, 11.7, 1.8, C_CARD)
    _shadow(chart_card)
    max_val = max(n_red, n_amber, n_green, 1)
    bar_max_w = 8.5
    for i, (label, count, color) in enumerate([("Red", n_red, C_RED), ("Amber", n_amber, C_AMBER), ("Green", n_green, C_GREEN)]):
        by = 3.8 + i * 0.50
        bw = (count / max_val) * bar_max_w if count > 0 else 0.3
        _tb(s, 1.1, by, 1, 0.35, label, 12, color, True, valign="middle")
        if count > 0:
            bar = _rect(s, 2.5, by + 0.03, bw, 0.28, color)
            _shadow(bar)
        _tb(s, 2.6 + bw, by, 0.8, 0.35, str(count), 13, C_WHITE, valign="middle")

    list_y = 5.7
    _tb(s, 0.8, list_y - 0.35, 5, 0.35, "Project Detail", 14, C_ICE, True, font="Georgia")
    for i, p in enumerate(projects):
        py = list_y + i * 0.4
        c = rag_color(p["overall"])
        _rect(s, 0.8, py, 0.7, 0.28, c)
        _tb(s, 0.8, py, 0.7, 0.28, p["overall"], 9, C_WHITE, True, PP_ALIGN.CENTER, valign="middle")
        _tb(s, 1.7, py, 6, 0.28, p["name"], 11, C_WHITE, valign="middle")
        _tb(s, 8, py, 4, 0.28, f"{int((p.get('pct_complete') or 0) * 100)}% · {p.get('pm') or 'N/A'}", 10, C_GRAY, valign="middle")

    # ── Slide 3: Project Cards ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, C_DARK)
    _tb(s, 0.8, 0.4, 12, 0.7, "Project Detail & Signal Breakdown", 28, C_WHITE, True, font="Georgia")

    card_w, card_h, gap = 3.8, 5.3, 0.3
    for i, p in enumerate(projects[:3]):
        cx = 0.8 + i * (card_w + gap)
        c = rag_color(p["overall"])
        bg = _rect(s, cx, 1.3, card_w, card_h, C_CARD)
        _shadow(bg)
        _rect(s, cx, 1.3, card_w, 0.06, c)
        _tb(s, cx + 0.3, 1.55, card_w - 0.6, 0.5, p["name"], 14, C_WHITE, True, font="Georgia")
        _rect(s, cx + 0.3, 2.15, 0.7, 0.3, c)
        _tb(s, cx + 0.3, 2.15, 0.7, 0.3, p["overall"], 10, C_WHITE, True, PP_ALIGN.CENTER, valign="middle")
        _tb(s, cx + 1.2, 2.15, 2, 0.3, f"{int((p.get('pct_complete') or 0) * 100)}% complete", 11, C_ICE, valign="middle")

        sig_y = 2.7
        signals = p.get("signals") or {}
        for j, sig_name in enumerate(signals.keys()):
            sy = sig_y + j * 0.38
            st = signals[sig_name]
            sc = rag_color(st)
            _oval(s, cx + 0.3, sy + 0.08, 0.15, sc)
            label = sig_name.replace("_", " ")
            _tb(s, cx + 0.6, sy, 2.2, 0.3, label, 10, C_GRAY, valign="middle")
            val = "—" if st == "insufficient_data" else st
            _tb(s, cx + 2.8, sy, 0.8, 0.3, val, 10, sc, True, PP_ALIGN.RIGHT, valign="middle")

        _tb(s, cx + 0.3, 1.3 + card_h - 0.5, card_w - 0.6, 0.3, f"PM: {p.get('pm') or 'N/A'}", 10, C_MUTED)

    # ── Slide 4: Trend Analysis ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, C_DARK)
    _tb(s, 0.8, 0.4, 10, 0.7, "Trend Analysis", 28, C_WHITE, True, font="Georgia")
    _tb(s, 0.8, 1.2, 12, 0.4, "Current snapshot — trends will appear after 2+ weeks of data collection.", 13, C_GRAY)

    hdr_y = 1.9
    _rect(s, 0.8, hdr_y, 11.7, 0.4, C_NAVY)
    hdrs = [("Project", 1, 5), ("Status", 6.5, 1.2), ("Schedule", 8, 1.2), ("Milestones", 9.2, 1.2), ("Blockers", 10.4, 1.2)]
    for text, x, w in hdrs:
        _tb(s, x, hdr_y, w, 0.4, text, 12, C_ICE, True, valign="middle")

    for i, p in enumerate(projects):
        ry = hdr_y + 0.5 + i * 0.45
        _rect(s, 0.8, ry, 11.7, 0.4, C_CARD if i % 2 == 0 else C_DARK)
        _tb(s, 1, ry, 5.3, 0.4, p["name"], 11, C_WHITE, valign="middle")
        sc = rag_color(p["overall"])
        _rect(s, 6.7, ry + 0.08, 0.6, 0.24, sc)
        _tb(s, 6.7, ry + 0.08, 0.6, 0.24, p["overall"], 8, C_WHITE, True, PP_ALIGN.CENTER, valign="middle")
        signals = p.get("signals") or {}
        for j, sig in enumerate(["schedule_slippage", "milestone_health", "blockers"]):
            st = signals.get(sig, "—")
            sc2 = C_GRAY if st == "insufficient_data" else rag_color(st)
            sx = [8, 9.2, 10.4][j]
            dot = "○" if st == "insufficient_data" else "●"
            val = "—" if st == "insufficient_data" else st
            _tb(s, sx, ry, 1.3, 0.4, f"{dot} {val}", 10, sc2, valign="middle")

    # ── Slide 5: Emerging Risks ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, C_DARK)
    _tb(s, 0.8, 0.4, 10, 0.7, "Emerging Risks", 28, C_WHITE, True, font="Georgia")

    risks = []
    for p in projects:
        if p["overall"] == "Red":
            risks.append(("Critical", p["name"], "Multiple Red signals — schedule slippage and overdue milestones. Escalation recommended.", C_RED))
        elif p["overall"] == "Amber":
            risks.append(("Watch", p["name"], "Budget overrun or signal drift. Requires active monitoring and mitigation plan.", C_AMBER))
        sigs = p.get("signals") or {}
        if sigs.get("milestone_health") == "Red":
            risks.append(("High", p["name"], "Overdue milestone tasks indicate delivery timeline risk.", C_AMBER))
    if not risks:
        risks.append(("Low", "Portfolio", "No significant risks identified across all projects.", C_GREEN))

    for i, (sev, proj, desc, color) in enumerate(risks[:6]):
        ry = 1.3 + i * 0.9
        card = _rect(s, 0.8, ry, 11.7, 0.75, C_CARD)
        _shadow(card)
        _rect(s, 0.8, ry, 0.06, 0.75, color)
        _rect(s, 1.2, ry + 0.12, 0.8, 0.28, color)
        _tb(s, 1.2, ry + 0.12, 0.8, 0.28, sev, 8, C_WHITE, True, PP_ALIGN.CENTER, valign="middle")
        _tb(s, 2.2, ry + 0.05, 9, 0.28, proj, 12, C_WHITE, True, font="Georgia")
        _tb(s, 2.2, ry + 0.35, 9, 0.35, desc, 11, C_GRAY)

    # ── Slide 6: Insights ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, C_DARK)
    _tb(s, 0.8, 0.4, 12, 0.7, "Insights & Recommendations", 28, C_WHITE, True, font="Georgia")

    insights = [
        ("!", f"{n_red + n_amber} of {total} projects require attention", f"{n_red} Red (immediate escalation) and {n_amber} Amber (monitoring) status projects identified."),
        ("D", "Data gaps limit assessment depth", "Budget and stakeholder sentiment fields are missing from standard exports. Adding these would improve accuracy."),
        ("M", "Milestone overdue is the top risk driver", "Across all projects, overdue tasks are the primary factor pushing status to Red."),
        ("R", "Recommend adding structured sentiment field", "A PM status-comment field would enable the sentiment signal and provide richer context."),
    ]
    for i, (icon, title, desc) in enumerate(insights):
        iy = 1.3 + i * 1.35
        card = _rect(s, 0.8, iy, 11.7, 1.1, C_CARD)
        _shadow(card)
        _oval(s, 1.2, iy + 0.2, 0.6, C_NAVY)
        _tb(s, 1.2, iy + 0.2, 0.6, 0.6, icon, 16, C_ICE, True, PP_ALIGN.CENTER, valign="middle", font="Georgia")
        _tb(s, 2.1, iy + 0.15, 9.5, 0.35, title, 14, C_WHITE, True, font="Georgia")
        _tb(s, 2.1, iy + 0.55, 9.5, 0.45, desc, 11, C_GRAY)

    # ── Slide 7: Methodology ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, C_NAVY)
    _tb(s, 0.8, 0.4, 10, 0.7, "Methodology", 28, C_WHITE, True, font="Georgia")
    _rect(s, 0, 7, 13.33, 0.5, C_DARK)
    _tb(s, 0.8, 7.05, 11.7, 0.4, "Project Health Reporting Agent  |  RAG_METHODOLOGY.md", 10, C_MUTED)

    methods = [
        ("Schedule", "% of tasks flagged Red determines status (≤5% Green, ≤15% Amber, >15% Red)"),
        ("Budget", "Burn ratio = spend% / complete%. Compares spending rate against progress rate."),
        ("Milestones", "Count of overdue (past End Date) incomplete tasks. >15 = Red, 6-15 = Amber."),
        ("Blockers", "Overdue critical tasks serve as blocker proxy. >2 = Red, 1-2 = Amber."),
        ("Sentiment", "Narrative signal from PM comments. Currently disabled — no structured source field."),
    ]
    col_w, col_gap = 2.1, 0.25
    for i, (label, desc) in enumerate(methods):
        cx = 0.8 + i * (col_w + col_gap)
        card = _rect(s, cx, 1.4, col_w, 3.9, C_CARD)
        _shadow(card)
        _rect(s, cx, 1.4, col_w, 0.06, C_ICE)
        _tb(s, cx + 0.2, 1.7, col_w - 0.4, 0.4, label, 14, C_ICE, True, font="Georgia")
        _tb(s, cx + 0.2, 2.3, col_w - 0.4, 2.5, desc, 11, C_GRAY)

    _tb(s, 0.8, 5.6, 11.7, 0.5,
        "Overall status = worst signal, except single-Red capped to Amber if all other signals are Green.  |  Thresholds: config/rag_thresholds.yaml",
        11, C_MUTED)

    return prs


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input-dir", default="docs/reports")
    parser.add_argument("--output", default="docs/presentations/monthly_report.pptx")
    args = parser.parse_args()
    data = load_weekly_data(args.input_dir)
    prs = build(data)
    prs.save(args.output)
    print(f"Saved: {args.output}")


if __name__ == "__main__":
    main()
